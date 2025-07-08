"""
Enhanced Google Drive Tool for comprehensive file processing.
Supports all major file types with complete content extraction.
"""

import os
import json
import io
import logging
from typing import Type, List, Dict, Any, Optional, Tuple
from pydantic import BaseModel, Field, PrivateAttr
from crewai.tools import BaseTool

from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from googleapiclient.errors import HttpError

# Configure logging
logger = logging.getLogger(__name__)

# Scopes for Google Drive API
SCOPES = [
    'https://www.googleapis.com/auth/drive.readonly',
    'https://www.googleapis.com/auth/documents.readonly',
    'https://www.googleapis.com/auth/spreadsheets.readonly',
    'https://www.googleapis.com/auth/presentations.readonly'
]

class GoogleDriveToolInput(BaseModel):
    """Input schema for GoogleDriveTool."""
    folder_id: str = Field(..., description="Google Drive folder ID to process for comprehensive analysis")

class GoogleDriveTool(BaseTool):
    name: str = "Google Drive Comprehensive File Processor"
    description: str = (
        "Advanced Google Drive tool that comprehensively processes ALL files in a specified folder. "
        "Handles PDF, DOCX, TXT, MD, CSV, Google Docs, Google Sheets, Google Slides, and more. "
        "Ensures complete content extraction from all pages and sections of each document for "
        "thorough requirements analysis. Processes files recursively and provides detailed progress tracking."
    )
    args_schema: Type[BaseModel] = GoogleDriveToolInput
    
    # Use private attributes for internal state
    _service: Optional[Any] = PrivateAttr(default=None)
    _processing_stats: Dict[str, Any] = PrivateAttr(default_factory=lambda: {
        'total_files': 0,
        'processed_successfully': 0,
        'failed_files': 0,
        'total_content_chars': 0,
        'file_types_processed': set()
    })

    def _load_credentials(self) -> Dict[str, Any]:
        """Load and validate credentials from credentials.json file."""
        credentials_path = 'credentials.json'
        
        if not os.path.exists(credentials_path):
            raise FileNotFoundError(
                f"Credentials file not found: {credentials_path}\n"
                "Please create credentials.json using credentials.json.example as template"
            )
        
        try:
            with open(credentials_path, 'r') as f:
                credentials_data = json.load(f)
            
            # Validate and normalize credential structure
            if 'installed' in credentials_data:
                return credentials_data
            elif 'web' in credentials_data:
                # Convert web credentials to installed format
                web_config = credentials_data['web']
                return {
                    'installed': {
                        'client_id': web_config['client_id'],
                        'project_id': web_config.get('project_id'),
                        'auth_uri': web_config.get('auth_uri', 'https://accounts.google.com/o/oauth2/auth'),
                        'token_uri': web_config.get('token_uri', 'https://oauth2.googleapis.com/token'),
                        'auth_provider_x509_cert_url': web_config.get(
                            'auth_provider_x509_cert_url', 
                            'https://www.googleapis.com/oauth2/v1/certs'
                        ),
                        'client_secret': web_config['client_secret'],
                        'redirect_uris': web_config.get('redirect_uris', ['http://localhost'])
                    }
                }
            else:
                raise ValueError(
                    "Invalid credentials format. Expected 'installed' or 'web' configuration.\n"
                    "Please check credentials.json.example for correct format."
                )
                
        except json.JSONDecodeError as e:
            raise ValueError(f"Invalid JSON in credentials file: {e}")
        except KeyError as e:
            raise ValueError(f"Missing required field in credentials: {e}")
        except Exception as e:
            raise Exception(f"Error reading credentials file: {e}")

    def _authenticate(self) -> Any:
        """Authenticate with Google Drive API using OAuth2."""
        logger.info("🔐 Authenticating with Google Drive API...")
        
        creds = None
        
        # Load existing token if available
        if os.path.exists('token.json'):
            try:
                creds = Credentials.from_authorized_user_file('token.json', SCOPES)
                logger.debug("Loaded existing token")
            except Exception as e:
                logger.warning(f"Failed to load existing token: {e}")
                # Remove corrupted token
                try:
                    os.remove('token.json')
                except:
                    pass
                creds = None
        
        # Refresh or re-authenticate if needed
        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                try:
                    logger.info("🔄 Refreshing expired token...")
                    creds.refresh(Request())
                    logger.info("✅ Token refreshed successfully")
                except Exception as e:
                    logger.warning(f"Token refresh failed: {e}")
                    creds = None
            
            if not creds:
                logger.info("🌐 Starting OAuth flow...")
                # Load credentials and start OAuth flow
                credentials_config = self._load_credentials()
                
                # Create temporary file for OAuth flow
                temp_creds_path = 'temp_credentials.json'
                try:
                    with open(temp_creds_path, 'w') as f:
                        json.dump(credentials_config, f)
                    
                    flow = InstalledAppFlow.from_client_secrets_file(temp_creds_path, SCOPES)
                    creds = flow.run_local_server(port=0, open_browser=True)
                    logger.info("✅ OAuth authentication successful")
                    
                finally:
                    # Clean up temporary file
                    if os.path.exists(temp_creds_path):
                        os.remove(temp_creds_path)
            
            # Save credentials for next run
            with open('token.json', 'w') as token:
                token.write(creds.to_json())
                logger.debug("Saved authentication token")
        
        # Build and return service
        service = build('drive', 'v3', credentials=creds)
        logger.info("✅ Google Drive API authentication complete")
        return service

    def _get_files_from_folder(self, service: Any, folder_id: str, recursive: bool = True) -> List[Dict]:
        """Get comprehensive list of files from folder and optionally subfolders."""
        logger.info(f"📁 Scanning folder: {folder_id}")
        
        try:
            # Verify folder exists and get info
            try:
                folder = service.files().get(fileId=folder_id).execute()
                folder_name = folder.get('name', 'Unknown')
                logger.info(f"📂 Processing folder: '{folder_name}'")
            except HttpError as e:
                if e.resp.status == 404:
                    return [{"error": f"Folder {folder_id} not found. Please check the folder ID and permissions."}]
                elif e.resp.status == 403:
                    return [{"error": f"Access denied to folder {folder_id}. Please check sharing permissions."}]
                else:
                    return [{"error": f"Cannot access folder {folder_id}: {str(e)}"}]
            except Exception as e:
                return [{"error": f"Cannot access folder {folder_id}: {str(e)}"}]
            
            all_files = []
            
            # Get files directly in this folder
            files = self._get_files_in_directory(service, folder_id)
            all_files.extend(files)
            
            # Recursively process subfolders if enabled
            if recursive:
                subfolders = self._get_subfolders(service, folder_id)
                for subfolder in subfolders:
                    subfolder_id = subfolder['id']
                    subfolder_name = subfolder['name']
                    logger.info(f"📁 Processing subfolder: '{subfolder_name}'")
                    
                    subfolder_files = self._get_files_in_directory(service, subfolder_id)
                    # Add subfolder context to files
                    for file_info in subfolder_files:
                        if 'error' not in file_info:
                            file_info['subfolder'] = subfolder_name
                    all_files.extend(subfolder_files)
            
            if not all_files:
                return [{"error": f"No files found in folder {folder_id} or its subfolders"}]
            
            # Filter out error entries for counting
            valid_files = [f for f in all_files if 'error' not in f]
            error_files = [f for f in all_files if 'error' in f]
            
            logger.info(f"✅ Found {len(valid_files)} files total")
            if error_files:
                logger.warning(f"⚠️ {len(error_files)} errors encountered during scanning")
            
            return all_files
            
        except Exception as e:
            error_msg = f"Failed to scan folder {folder_id}: {str(e)}"
            logger.error(error_msg)
            return [{"error": error_msg}]

    def _get_files_in_directory(self, service: Any, directory_id: str) -> List[Dict]:
        """Get files from a specific directory."""
        try:
            query = f"'{directory_id}' in parents and trashed=false"
            
            files = []
            page_token = None
            
            while True:
                results = service.files().list(
                    q=query,
                    fields="nextPageToken, files(id, name, mimeType, size, modifiedTime, parents)",
                    pageSize=100,
                    pageToken=page_token
                ).execute()
                
                items = results.get('files', [])
                
                # Filter out folders from this list (we handle them separately)
                for item in items:
                    if item.get('mimeType') != 'application/vnd.google-apps.folder':
                        files.append(item)
                
                page_token = results.get('nextPageToken')
                if not page_token:
                    break
            
            return files
            
        except Exception as e:
            return [{"error": f"Failed to get files from directory {directory_id}: {str(e)}"}]

    def _get_subfolders(self, service: Any, parent_folder_id: str) -> List[Dict]:
        """Get list of subfolders in a parent folder."""
        try:
            query = f"'{parent_folder_id}' in parents and mimeType='application/vnd.google-apps.folder' and trashed=false"
            
            results = service.files().list(
                q=query,
                fields="files(id, name)",
                pageSize=100
            ).execute()
            
            return results.get('files', [])
            
        except Exception as e:
            logger.warning(f"Failed to get subfolders from {parent_folder_id}: {str(e)}")
            return []

    def _extract_content_by_type(self, service: Any, file_id: str, file_name: str, mime_type: str) -> Tuple[str, bool]:
        """Extract content based on file type. Returns (content, success)."""
        try:
            logger.debug(f"🔄 Processing: {file_name} ({mime_type})")
            
            # Handle Google Workspace files
            if 'google-apps.document' in mime_type:
                return self._extract_google_doc_content(service, file_id), True
            elif 'google-apps.spreadsheet' in mime_type:
                return self._extract_google_sheet_content(service, file_id), True
            elif 'google-apps.presentation' in mime_type:
                return self._extract_google_slides_content(service, file_id), True
            
            # Handle regular files by extension and MIME type
            file_ext = file_name.lower().split('.')[-1] if '.' in file_name else ''
            
            if file_ext == 'pdf' or 'pdf' in mime_type:
                return self._extract_pdf_content(service, file_id), True
            elif file_ext in ['txt', 'md', 'csv', 'py', 'js', 'html', 'css', 'json', 'xml'] or 'text/' in mime_type:
                return self._extract_text_content(service, file_id), True
            elif file_ext in ['docx', 'doc'] or 'word' in mime_type or 'document' in mime_type:
                return self._extract_docx_content(service, file_id), True
            elif file_ext in ['xlsx', 'xls'] or 'spreadsheet' in mime_type:
                return self._extract_excel_content(service, file_id), True
            elif file_ext in ['pptx', 'ppt'] or 'presentation' in mime_type:
                return self._extract_powerpoint_content(service, file_id), True
            else:
                # Try to extract as text for unknown types
                try:
                    content = self._extract_text_content(service, file_id)
                    if content and len(content.strip()) > 0:
                        return content, True
                except:
                    pass
                
                return f"Unsupported file type: {file_name} ({mime_type})", False
                
        except Exception as e:
            error_msg = f"Error extracting content from {file_name}: {str(e)}"
            logger.error(error_msg)
            return error_msg, False

    def _extract_text_content(self, service: Any, file_id: str) -> str:
        """Extract content from text-based files."""
        request = service.files().get_media(fileId=file_id)
        file_io = io.BytesIO()
        downloader = MediaIoBaseDownload(file_io, request)
        
        done = False
        while not done:
            status, done = downloader.next_chunk()
        
        # Try multiple encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                content = file_io.getvalue().decode(encoding)
                return content
            except UnicodeDecodeError:
                continue
        
        # Fallback: decode with error handling
        return file_io.getvalue().decode('utf-8', errors='replace')

    def _extract_pdf_content(self, service: Any, file_id: str) -> str:
        """Extract comprehensive content from PDF files."""
        try:
            import PyPDF2
        except ImportError:
            return "Error: PyPDF2 not installed. Install with: pip install PyPDF2"
        
        request = service.files().get_media(fileId=file_id)
        file_io = io.BytesIO()
        downloader = MediaIoBaseDownload(file_io, request)
        
        done = False
        while not done:
            status, done = downloader.next_chunk()
        
        file_io.seek(0)
        try:
            pdf_reader = PyPDF2.PdfReader(file_io)
            content_parts = []
            
            # Add metadata
            if pdf_reader.metadata:
                content_parts.append("=== PDF METADATA ===")
                for key, value in pdf_reader.metadata.items():
                    content_parts.append(f"{key}: {value}")
                content_parts.append("\n=== CONTENT ===")
            
            # Extract text from all pages
            total_pages = len(pdf_reader.pages)
            logger.debug(f"📄 Extracting from {total_pages} PDF pages")
            
            for page_num, page in enumerate(pdf_reader.pages, 1):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        content_parts.append(f"\n--- PAGE {page_num} of {total_pages} ---")
                        content_parts.append(page_text.strip())
                except Exception as e:
                    content_parts.append(f"\n--- PAGE {page_num} (Error: {str(e)}) ---")
            
            return "\n".join(content_parts)
        except Exception as e:
            return f"Error reading PDF: {str(e)}"

    def _extract_docx_content(self, service: Any, file_id: str) -> str:
        """Extract comprehensive content from DOCX files."""
        try:
            import docx
        except ImportError:
            return "Error: python-docx not installed. Install with: pip install python-docx"
        
        request = service.files().get_media(fileId=file_id)
        file_io = io.BytesIO()
        downloader = MediaIoBaseDownload(file_io, request)
        
        done = False
        while not done:
            status, done = downloader.next_chunk()
        
        file_io.seek(0)
        try:
            doc = docx.Document(file_io)
            content_parts = []
            
            # Extract paragraphs
            content_parts.append("=== DOCUMENT CONTENT ===")
            for para in doc.paragraphs:
                if para.text.strip():
                    content_parts.append(para.text.strip())
            
            # Extract tables
            if doc.tables:
                content_parts.append("\n=== TABLES ===")
                for table_num, table in enumerate(doc.tables, 1):
                    content_parts.append(f"\n--- TABLE {table_num} ---")
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        content_parts.append(" | ".join(row_data))
            
            return "\n".join(content_parts)
        except Exception as e:
            return f"Error reading DOCX: {str(e)}"

    def _extract_excel_content(self, service: Any, file_id: str) -> str:
        """Extract content from Excel files."""
        try:
            import openpyxl
        except ImportError:
            return "Error: openpyxl not installed. Install with: pip install openpyxl"
        
        request = service.files().get_media(fileId=file_id)
        file_io = io.BytesIO()
        downloader = MediaIoBaseDownload(file_io, request)
        
        done = False
        while not done:
            status, done = downloader.next_chunk()
        
        file_io.seek(0)
        try:
            workbook = openpyxl.load_workbook(file_io, data_only=True)
            content_parts = []
            
            content_parts.append("=== EXCEL WORKBOOK ===")
            
            for sheet_name in workbook.sheetnames:
                worksheet = workbook[sheet_name]
                content_parts.append(f"\n--- SHEET: {sheet_name} ---")
                
                # Get all data from sheet
                for row in worksheet.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    if any(cell.strip() for cell in row_data):  # Skip empty rows
                        content_parts.append(" | ".join(row_data))
            
            return "\n".join(content_parts)
        except Exception as e:
            return f"Error reading Excel file: {str(e)}"

    def _extract_powerpoint_content(self, service: Any, file_id: str) -> str:
        """Extract content from PowerPoint files."""
        try:
            from pptx import Presentation
        except ImportError:
            return "Error: python-pptx not installed. Install with: pip install python-pptx"
        
        request = service.files().get_media(fileId=file_id)
        file_io = io.BytesIO()
        downloader = MediaIoBaseDownload(file_io, request)
        
        done = False
        while not done:
            status, done = downloader.next_chunk()
        
        file_io.seek(0)
        try:
            prs = Presentation(file_io)
            content_parts = []
            
            content_parts.append("=== POWERPOINT PRESENTATION ===")
            
            for slide_num, slide in enumerate(prs.slides, 1):
                content_parts.append(f"\n--- SLIDE {slide_num} ---")
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        content_parts.append(shape.text.strip())
                    
                    # Extract table content if present
                    if hasattr(shape, 'table'):
                        content_parts.append("\nTable:")
                        for row in shape.table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            content_parts.append(" | ".join(row_data))
            
            return "\n".join(content_parts)
        except Exception as e:
            return f"Error reading PowerPoint file: {str(e)}"

    def _extract_google_doc_content(self, service: Any, file_id: str) -> str:
        """Extract comprehensive content from Google Docs."""
        try:
            # Export as plain text for main content
            request = service.files().export_media(fileId=file_id, mimeType='text/plain')
            file_io = io.BytesIO()
            downloader = MediaIoBaseDownload(file_io, request)
            
            done = False
            while not done:
                status, done = downloader.next_chunk()
            
            content = file_io.getvalue().decode('utf-8')
            
            # Also try to get rich text version for additional formatting info
            try:
                request_rich = service.files().export_media(fileId=file_id, mimeType='text/html')
                file_io_rich = io.BytesIO()
                downloader_rich = MediaIoBaseDownload(file_io_rich, request_rich)
                
                done = False
                while not done:
                    status, done = downloader_rich.next_chunk()
                
                html_content = file_io_rich.getvalue().decode('utf-8')
                
                # Extract additional text from HTML if needed
                try:
                    from bs4 import BeautifulSoup
                    soup = BeautifulSoup(html_content, 'html.parser')
                    # Remove script and style elements
                    for script in soup(["script", "style"]):
                        script.decompose()
                    html_text = soup.get_text()
                    
                    # If HTML version has significantly more content, use it
                    if len(html_text.strip()) > len(content.strip()) * 1.1:
                        content = html_text
                except ImportError:
                    pass  # BeautifulSoup not available, use plain text
                except Exception:
                    pass  # HTML parsing failed, use plain text
            except:
                pass  # Rich text export failed, use plain text
            
            return content
        except Exception as e:
            return f"Error extracting Google Doc: {str(e)}"

    def _extract_google_sheet_content(self, service: Any, file_id: str) -> str:
        """Extract comprehensive content from Google Sheets."""
        try:
            # Export as CSV for structured data
            request = service.files().export_media(fileId=file_id, mimeType='text/csv')
            file_io = io.BytesIO()
            downloader = MediaIoBaseDownload(file_io, request)
            
            done = False
            while not done:
                status, done = downloader.next_chunk()
            
            csv_content = file_io.getvalue().decode('utf-8')
            
            # Also try to get Excel format for multiple sheets
            try:
                request_excel = service.files().export_media(
                    fileId=file_id, 
                    mimeType='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
                )
                file_io_excel = io.BytesIO()
                downloader_excel = MediaIoBaseDownload(file_io_excel, request_excel)
                
                done = False
                while not done:
                    status, done = downloader_excel.next_chunk()
                
                # Process Excel file to get all sheets
                try:
                    import openpyxl
                    file_io_excel.seek(0)
                    workbook = openpyxl.load_workbook(file_io_excel, data_only=True)
                    
                    content_parts = ["=== GOOGLE SHEETS (ALL SHEETS) ==="]
                    
                    for sheet_name in workbook.sheetnames:
                        worksheet = workbook[sheet_name]
                        content_parts.append(f"\n--- SHEET: {sheet_name} ---")
                        
                        for row in worksheet.iter_rows(values_only=True):
                            row_data = [str(cell) if cell is not None else "" for cell in row]
                            if any(cell.strip() for cell in row_data):
                                content_parts.append(" | ".join(row_data))
                    
                    return "\n".join(content_parts)
                except:
                    pass  # Fall back to CSV
            except:
                pass  # Excel export failed, use CSV
            
            return f"=== GOOGLE SHEETS ===\n{csv_content}"
        except Exception as e:
            return f"Error extracting Google Sheet: {str(e)}"

    def _extract_google_slides_content(self, service: Any, file_id: str) -> str:
        """Extract content from Google Slides."""
        try:
            # Export as plain text
            request = service.files().export_media(fileId=file_id, mimeType='text/plain')
            file_io = io.BytesIO()
            downloader = MediaIoBaseDownload(file_io, request)
            
            done = False
            while not done:
                status, done = downloader.next_chunk()
            
            content = file_io.getvalue().decode('utf-8')
            return f"=== GOOGLE SLIDES ===\n{content}"
        except Exception as e:
            return f"Error extracting Google Slides: {str(e)}"

    def _update_stats(self, file_name: str, mime_type: str, content_length: int, success: bool):
        """Update processing statistics."""
        stats = self._processing_stats
        stats['total_files'] += 1
        if success:
            stats['processed_successfully'] += 1
            stats['total_content_chars'] += content_length
        else:
            stats['failed_files'] += 1
        
        # Track file types
        file_ext = file_name.split('.')[-1].lower() if '.' in file_name else 'unknown'
        stats['file_types_processed'].add(file_ext)

    def _generate_summary(self) -> str:
        """Generate processing summary."""
        stats = self._processing_stats
        
        summary_parts = [
            "\n" + "="*80,
            "📊 PROCESSING SUMMARY",
            "="*80,
            f"📁 Total files found: {stats['total_files']}",
            f"✅ Successfully processed: {stats['processed_successfully']}",
            f"❌ Failed to process: {stats['failed_files']}",
            f"📝 Total content extracted: {stats['total_content_chars']:,} characters",
            f"📄 File types processed: {', '.join(sorted(stats['file_types_processed']))}",
            f"📈 Success rate: {(stats['processed_successfully']/max(stats['total_files'], 1)*100):.1f}%",
            "="*80
        ]
        
        return "\n".join(summary_parts)

    def _run(self, folder_id: str) -> str:
        """Main execution method for comprehensive file processing."""
        logger.info(f"🚀 Starting comprehensive Google Drive analysis")
        logger.info(f"📁 Target folder ID: {folder_id}")
        
        try:
            # Reset stats
            self._processing_stats = {
                'total_files': 0,
                'processed_successfully': 0,
                'failed_files': 0,
                'total_content_chars': 0,
                'file_types_processed': set()
            }
            
            # Authenticate
            service = self._authenticate()
            
            # Get all files from folder
            files = self._get_files_from_folder(service, folder_id, recursive=True)
            
            if not files or (len(files) == 1 and "error" in files[0]):
                error_msg = files[0].get("error", "No files found") if files else "No files found"
                logger.error(f"❌ {error_msg}")
                return error_msg
            
            # Process each file
            all_content = []
            
            for file_info in files:
                if "error" in file_info:
                    logger.warning(f"⚠️ Skipping file due to error: {file_info['error']}")
                    continue
                
                file_name = file_info['name']
                file_id = file_info['id']
                mime_type = file_info.get('mimeType', '')
                subfolder = file_info.get('subfolder', '')
                
                logger.info(f"🔄 Processing: {file_name}")
                
                content, success = self._extract_content_by_type(service, file_id, file_name, mime_type)
                
                # Update statistics
                self._update_stats(file_name, mime_type, len(content), success)
                
                # Add content to results
                if success and content and not content.startswith("Error") and not content.startswith("Unsupported"):
                    all_content.append(f"\n{'='*80}")
                    all_content.append(f"FILE: {file_name}")
                    if subfolder:
                        all_content.append(f"SUBFOLDER: {subfolder}")
                    all_content.append(f"TYPE: {mime_type}")
                    all_content.append(f"SIZE: {file_info.get('size', 'Unknown')} bytes")
                    all_content.append(f"MODIFIED: {file_info.get('modifiedTime', 'Unknown')}")
                    all_content.append(f"{'='*80}")
                    all_content.append(content)
                    logger.info(f"✅ Successfully processed: {file_name} ({len(content):,} chars)")
                else:
                    logger.warning(f"⚠️ Could not process: {file_name} - {content[:100]}...")
                    all_content.append(f"\n❌ Failed to process {file_name}: {content}")
            
            # Generate final content with summary
            final_content = "\n".join(all_content)
            summary = self._generate_summary()
            
            result = summary + "\n" + final_content
            
            logger.info(f"🎉 Processing complete! {self._processing_stats['processed_successfully']} files processed successfully")
            
            return result
            
        except Exception as e:
            error_msg = f"💥 Error processing Google Drive folder: {str(e)}"
            logger.error(error_msg)
            return error_msg